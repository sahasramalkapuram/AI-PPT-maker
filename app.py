import os
import json
import io
from flask import Flask, render_template, request, redirect, url_for, send_file
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

app = Flask(__name__, template_folder='templates')
app.secret_key = os.environ.get("SECRET_KEY", "dev-secret-key-12345")

# Configure Gemini AI using the environment variable key
genai.configure(api_key=os.environ.get("GEMINI_API_KEY"))

def generate_real_ai_slides(prompt, count, assets_config):
    """
    Calls Gemini AI to generate real, structured content for the slides.
    """
    try:
        model = genai.GenerativeModel('gemini-2.5-flash')
        
        # Explicit prompt forcing Gemini to return a clean JSON array matching our UI
        ai_prompt = f"""
        Create a presentation outline about "{prompt}" with exactly {count} slides.
        Return the response as a raw JSON array of objects. Do not include markdown code blocks like ```json.
        Each object in the array MUST have these exact keys:
        - "title": A short, engaging slide title.
        - "bullets": An array of 3 detailed, informative bullet points explaining the slide topic.
        - "image_query": A simple 2-3 word keyword phrase for finding an image related to this slide.
        - "chart_data": True if this slide can use a data graph/chart layout, otherwise False.

        Example format:
        [
          {{"title": "Introduction to Parallel Worlds", "bullets": ["Bullet 1", "Bullet 2", "Bullet 3"], "image_query": "multiverse physics", "chart_data": false}}
        ]
        """
        
        response = model.generate_content(ai_prompt)
        text_data = response.text.strip()
        
        # Clean up any accidental markdown blocks if the model includes them
        if text_data.startswith("```"):
            text_data = text_data.split("```")[1]
            if text_data.startswith("json"):
                text_data = text_data[4:]
        
        slides = json.loads(text_data.strip())
        
        # Ensure chart_data matches formatting expected by frontend template
        for slide in slides:
            if slide.get("chart_data") and assets_config in ['both', 'charts']:
                slide["chart_data"] = [20, 40, 60, 80]
            else:
                slide["chart_data"] = None
                
            if assets_config not in ['both', 'images']:
                slide["image_query"] = None
                
        return slides
    except Exception as e:
        print(f"Gemini API Error: {str(e)}")
        # Fallback slides if API fails or key is missing
        return [
            {
                "title": f"Slide {i+1}: Overview of {prompt}",
                "bullets": ["Detailed AI generation requires an active GEMINI_API_KEY.", "Ensure environment variables are set up on Render.", "This is a placeholder slide."],
                "image_query": "space galaxy" if assets_config in ['both', 'images'] else None,
                "chart_data": [10, 20, 30, 40] if assets_config in ['both', 'charts'] else None
            } for i in range(int(count))
        ]

@app.route('/')
def home():
    return redirect(url_for('dashboard'))

@app.route('/dashboard')
def dashboard():
    return render_template('dashboard.html')

@app.route('/generate/outline', methods=['POST'])
def generate_outline():
    try:
        prompt = request.form.get('prompt', '').strip()
        aesthetic = request.form.get('aesthetic', 'cyberpunk')
        slide_count = request.form.get('slide_count', '7')
        assets = request.form.get('assets', 'both')
        
        if not prompt:
            return redirect(url_for('dashboard'))
            
        # Call the real Gemini AI generator function
        generated_deck = generate_real_ai_slides(prompt, slide_count, assets)
        
        return render_template(
            'outline.html', 
            slides=generated_deck, 
            aesthetic=aesthetic,
            prompt=prompt
        )
    except Exception as e:
        return f"Internal Processing Error: {str(e)}", 500

@app.route('/export/pptx', methods=['POST'])
def export_pptx():
    try:
        slides_data_raw = request.form.get('slides_json', '[]')
        slides_list = json.loads(slides_data_raw)
        
        prs = Presentation()
        
        for slide_item in slides_list:
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Dark Theme Canvas Style
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(15, 23, 42)
            
            # Slide Title Text Box
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(1.0))
            tf_title = title_box.text_frame
            p_title = tf_title.paragraphs[0]
            p_title.text = slide_item.get('title', 'Presentation Slide')
            p_title.font.size = Pt(28)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(255, 255, 255)
            
            # Content Bullets Text Box
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            
            for index, bullet_text in enumerate(slide_item.get('bullets', [])):
                p_bullet = tf_content.add_paragraph() if index > 0 else tf_content.paragraphs[0]
                p_bullet.text = f"• {bullet_text}"
                p_bullet.font.size = Pt(14)
                p_bullet.font.color.rgb = RGBColor(203, 213, 225)
                p_bullet.space_after = Pt(12)

            # Sidebar Asset Layout Box Indicator
            if slide_item.get('chart_data') or slide_item.get('image_query'):
                asset_box = slide.shapes.add_textbox(Inches(6.5), Inches(2.0), Inches(3.0), Inches(4.0))
                tf_asset = asset_box.text_frame
                tf_asset.word_wrap = True
                p_asset = tf_asset.paragraphs[0]
                p_asset.text = f"📊 [Visual Asset Slot]\nQuery: {slide_item.get('image_query', 'N/A')}"
                p_asset.font.size = Pt(12)
                p_asset.font.italic = True
                p_asset.font.color.rgb = RGBColor(129, 140, 248)
                
        binary_output = io.BytesIO()
        prs.save(binary_output)
        binary_output.seek(0)
        
        return send_file(
            binary_output,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name="AI_Generated_Presentation.pptx"
        )
    except Exception as e:
        return f"PPTX Engine Export Error: {str(e)}", 500

if __name__ == '__main__':
    app.run(debug=True)
